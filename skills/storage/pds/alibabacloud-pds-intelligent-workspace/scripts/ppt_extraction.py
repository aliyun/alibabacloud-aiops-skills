#!/usr/bin/env python3
"""
PDS Video Analysis PPT Extraction Script

Extracts PPT slides from video analysis results and generates a PPTX file.
Supports adding notes (page number, timestamp, etc.).
"""

import json
import requests
import argparse
from pathlib import Path
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Missing dependency: python-pptx")
    print("Run: pip install -r scripts/requirements.txt")
    exit(1)


def ms_to_timestamp(ms):
    """Convert milliseconds to timestamp format"""
    seconds = ms // 1000
    hours = seconds // 3600
    minutes = (seconds % 3600) // 60
    secs = seconds % 60
    return f"{hours:02d}:{minutes:02d}:{secs:02d}"


def download_image(url):
    """Download image into memory"""
    response = requests.get(url, timeout=30)
    response.raise_for_status()
    return BytesIO(response.content)


def create_pptx_from_video_analysis(result_json, output_path="output.pptx", keep_aspect_ratio=False):
    """
    Create a PPTX file from video analysis results.

    Args:
        result_json: Complete API response (dict or JSON file path)
        output_path: Output PPTX file path
        keep_aspect_ratio: Whether to maintain aspect ratio (default: False, fills entire slide)

    Returns:
        bool: True on success, False on failure
    """
    if isinstance(result_json, str):
        with open(result_json, 'r', encoding='utf-8') as f:
            result = json.load(f)
    else:
        result = result_json

    if "ppt_details" not in result or not result["ppt_details"]:
        print("No PPT content detected in this video")
        return False

    ppt_details_url = result["ppt_details"][0]
    print(f"Downloading PPT details: {ppt_details_url}")
    ppt_data = requests.get(ppt_details_url, timeout=30).json()

    ppt_data.sort(key=lambda x: x["PPTShotIndex"])
    print(f"Detected {len(ppt_data)} PPT pages")

    images_map = result.get("images", {})

    prs = Presentation()

    # 16:9 widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for i, ppt_page in enumerate(ppt_data, 1):
        image_path = ppt_page["ImagePath"]
        ppt_index = ppt_page["PPTShotIndex"]
        start_time = ms_to_timestamp(ppt_page["StartTime"])

        print(f"  - Processing page {i}/{len(ppt_data)} (index: {ppt_index}, time: {start_time})")

        if image_path not in images_map:
            print(f"    Warning: image path {image_path} not found in images, skipping")
            continue

        image_info = images_map[image_path]
        image_url = image_info.get("Url") or image_info.get("url")

        try:
            image_stream = download_image(image_url)
        except Exception as e:
            print(f"    Failed to download image: {e}")
            continue

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        if keep_aspect_ratio:
            add_picture_with_aspect_ratio(
                slide,
                image_stream,
                prs.slide_width,
                prs.slide_height
            )
        else:
            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height

            slide.shapes.add_picture(
                image_stream,
                left, top,
                width=width,
                height=height
            )

        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = (
            f"Page: {i}\n"
            f"Index: {ppt_index}\n"
            f"Appears at: {start_time}\n"
            f"Image path: {image_path}"
        )

    prs.save(output_path)
    print(f"PPTX file generated: {output_path}")
    print(f"   Total pages: {len(prs.slides)}")

    return True


def add_picture_with_aspect_ratio(slide, image_stream, slide_width, slide_height):
    """
    Insert an image while maintaining aspect ratio.

    Args:
        slide: PPTX slide object
        image_stream: Image stream (BytesIO)
        slide_width: Slide width
        slide_height: Slide height
    """
    from PIL import Image

    img = Image.open(image_stream)
    img_width, img_height = img.size
    img_aspect = img_width / img_height

    slide_aspect = slide_width / slide_height

    if img_aspect > slide_aspect:
        width = slide_width
        height = slide_width / img_aspect
        left = Inches(0)
        top = (slide_height - height) / 2
    else:
        height = slide_height
        width = slide_height * img_aspect
        top = Inches(0)
        left = (slide_width - width) / 2

    image_stream.seek(0)

    slide.shapes.add_picture(image_stream, left, top, width=width, height=height)


def validate_pptx(pptx_path, expected_slide_count):
    """
    Validate the generated PPTX file.

    Args:
        pptx_path: PPTX file path
        expected_slide_count: Expected number of slides

    Returns:
        bool: Whether validation passed
    """
    try:
        prs = Presentation(pptx_path)
        actual_count = len(prs.slides)

        print(f"\nPPTX Validation Results:")
        print(f"   File path: {pptx_path}")
        print(f"   Expected pages: {expected_slide_count}")
        print(f"   Actual pages: {actual_count}")

        if actual_count == expected_slide_count:
            print("   Page count matches")
        else:
            print("   Page count mismatch")

        missing_images = []
        for i, slide in enumerate(prs.slides, 1):
            has_picture = any(
                shape.shape_type == 13
                for shape in slide.shapes
            )
            if not has_picture:
                missing_images.append(i)

        if missing_images:
            print(f"   Pages missing images: {missing_images}")
        else:
            print("   All pages contain images")

        return actual_count == expected_slide_count and not missing_images

    except Exception as e:
        print(f"Validation failed: {e}")
        return False


def main():
    parser = argparse.ArgumentParser(
        description='PDS Video Analysis PPT Extraction Tool',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python ppt_extraction.py video_analysis_result.json
  python ppt_extraction.py video_analysis_result.json -o extracted.pptx
  python ppt_extraction.py video_analysis_result.json --keep-aspect-ratio
  python ppt_extraction.py video_analysis_result.json --validate
        """
    )

    parser.add_argument(
        'input_file',
        help='JSON result file path from the video analysis API'
    )

    parser.add_argument(
        '-o', '--output',
        default='extracted_ppt.pptx',
        help='Output PPTX file path (default: extracted_ppt.pptx)'
    )

    parser.add_argument(
        '--keep-aspect-ratio',
        action='store_true',
        help='Maintain image aspect ratio (default: fill entire slide)'
    )

    parser.add_argument(
        '--validate',
        action='store_true',
        help='Validate the PPTX file after generation'
    )

    args = parser.parse_args()

    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"File not found: {args.input_file}")
        return 1

    try:
        success = create_pptx_from_video_analysis(
            args.input_file,
            args.output,
            args.keep_aspect_ratio
        )

        if success and args.validate:
            with open(args.input_file, 'r', encoding='utf-8') as f:
                result = json.load(f)

            if "ppt_details" in result and result["ppt_details"]:
                ppt_details_url = result["ppt_details"][0]
                ppt_data = requests.get(ppt_details_url, timeout=30).json()
                expected_count = len(ppt_data)
                validate_pptx(args.output, expected_count)

        return 0 if success else 1

    except Exception as e:
        print(f"Processing failed: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    exit(main())
